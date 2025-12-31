"""Ingestion worker loop with leasing and stage updates."""
from __future__ import annotations

import csv
import io
import importlib.util
import json
import logging
import os
import re
import time
import subprocess
import statistics
from dataclasses import dataclass
from datetime import date, datetime, timedelta, timezone
from hashlib import sha256
from pathlib import Path
import shutil
from typing import Callable
from uuid import uuid4

from sqlalchemy import and_, or_, text

from PIL import Image
import pdfplumber
from pdfminer.high_level import extract_pages
from pdfminer.layout import LAParams, LTAnno, LTChar, LTTextContainer, LTTextLine
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
from python_calamine import CalamineWorkbook, CalamineError, ZipError, XmlError
from tabulate import tabulate

from api.db.session import SessionLocal
from api.models.file_ingestion import FileProcessingJob, IngestedFile, FileDerivative
from api.config import settings
from api.services.storage.service import get_storage_backend


LEASE_SECONDS = 180
SLEEP_SECONDS = 2
MAX_ATTEMPTS = 3
BACKOFF_BASE_SECONDS = 2
BACKOFF_MAX_SECONDS = 60
PIPELINE_STAGES = [
    "validating",
    "converting",
    "extracting",
    "ai_md",
    "thumb",
    "finalizing",
]

logger = logging.getLogger("ingestion.worker")
logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
_COMMON_SHORT_WORDS = {
    "a",
    "i",
    "an",
    "and",
    "as",
    "at",
    "be",
    "by",
    "can",
    "do",
    "for",
    "from",
    "had",
    "has",
    "have",
    "he",
    "if",
    "in",
    "into",
    "is",
    "it",
    "me",
    "might",
    "may",
    "no",
    "not",
    "of",
    "on",
    "or",
    "our",
    "out",
    "so",
    "than",
    "that",
    "the",
    "them",
    "then",
    "there",
    "these",
    "they",
    "this",
    "to",
    "up",
    "us",
    "was",
    "were",
    "we",
    "with",
    "will",
    "would",
    "you",
    "your",
}

_audio_transcriber: Callable[..., dict] | None = None
_youtube_transcriber: Callable[..., dict] | None = None


def _now() -> datetime:
    return datetime.now(timezone.utc)


@dataclass(frozen=True)
class DerivativePayload:
    kind: str
    storage_key: str
    mime: str
    size_bytes: int
    sha256: str | None
    content: bytes


class IngestionError(Exception):
    def __init__(self, code: str, message: str, retryable: bool = False):
        super().__init__(message)
        self.code = code
        self.retryable = retryable


def _claim_job(db, worker_id: str) -> FileProcessingJob | None:
    job = (
        db.query(FileProcessingJob)
        .filter(
            and_(
                FileProcessingJob.status == "queued",
                or_(
                    FileProcessingJob.lease_expires_at.is_(None),
                    FileProcessingJob.lease_expires_at < _now(),
                ),
            )
        )
        .order_by(FileProcessingJob.updated_at.asc())
        .with_for_update(skip_locked=True)
        .first()
    )
    if not job:
        return None

    job.status = "processing"
    job.stage = "validating"
    job.worker_id = worker_id
    job.lease_expires_at = _now() + timedelta(seconds=LEASE_SECONDS)
    job.started_at = job.started_at or _now()
    job.updated_at = _now()
    db.commit()
    logger.info("Claimed job %s (file_id=%s)", job.id, job.file_id)
    return job


def _refresh_lease(db, job: FileProcessingJob) -> None:
    job.lease_expires_at = _now() + timedelta(seconds=LEASE_SECONDS)
    job.updated_at = _now()
    db.commit()


def _set_stage(db, job: FileProcessingJob, stage: str) -> None:
    job.stage = stage
    job.updated_at = _now()
    db.commit()


def _mark_ready(db, job: FileProcessingJob) -> None:
    job.status = "ready"
    job.stage = "ready"
    job.finished_at = _now()
    job.updated_at = _now()
    job.worker_id = None
    job.lease_expires_at = None
    db.commit()
    logger.info("Marked ready file_id=%s", job.file_id)


def _compute_backoff_seconds(attempts: int) -> int:
    if attempts <= 0:
        return 0
    delay = BACKOFF_BASE_SECONDS * (2 ** (attempts - 1))
    return min(delay, BACKOFF_MAX_SECONDS)


def _retry_or_fail(db, job: FileProcessingJob, error: IngestionError) -> None:
    job.attempts = (job.attempts or 0) + 1
    job.error_code = error.code
    job.error_message = str(error)
    job.updated_at = _now()
    job.worker_id = None
    job.lease_expires_at = None

    if error.retryable and job.attempts < MAX_ATTEMPTS:
        job.status = "queued"
        job.stage = "queued"
        job.lease_expires_at = _now() + timedelta(seconds=_compute_backoff_seconds(job.attempts))
        db.commit()
        logger.warning(
            "Retrying job file_id=%s attempt=%s code=%s",
            job.file_id,
            job.attempts,
            error.code,
        )
        return

    job.status = "failed"
    job.stage = "failed"
    job.finished_at = _now()
    db.commit()
    logger.error(
        "Job failed file_id=%s attempt=%s code=%s message=%s",
        job.file_id,
        job.attempts,
        error.code,
        str(error),
    )


def _should_cleanup_after_failure(error: IngestionError) -> bool:
    if error.retryable:
        return False
    if error.code in {"CONVERSION_UNAVAILABLE"}:
        return False
    return True


def _requeue_stalled_jobs(db) -> None:
    stalled_jobs = (
        db.query(FileProcessingJob)
        .filter(
            and_(
                FileProcessingJob.status == "processing",
                FileProcessingJob.lease_expires_at.is_not(None),
                FileProcessingJob.lease_expires_at < _now(),
            )
        )
        .order_by(FileProcessingJob.updated_at.asc())
        .limit(5)
        .all()
    )
    for job in stalled_jobs:
        logger.warning("Requeuing stalled job file_id=%s", job.file_id)
        _retry_or_fail(
            db,
            job,
            IngestionError("WORKER_STALLED", "Worker heartbeat expired", retryable=True),
        )


def _get_file(db, job: FileProcessingJob) -> IngestedFile:
    record = db.query(IngestedFile).filter(IngestedFile.id == job.file_id).first()
    if not record:
        raise IngestionError("FILE_NOT_FOUND", "Ingestion record missing", retryable=False)
    return record


def _staging_path(file_id: str) -> Path:
    return Path("/tmp/sidebar-ingestion") / file_id / "source"


def _cleanup_staging(file_id: str) -> None:
    staging_root = Path("/tmp/sidebar-ingestion") / file_id
    if staging_root.exists():
        shutil.rmtree(staging_root, ignore_errors=True)


def _derivative_dir(file_id: str) -> Path:
    return Path("/tmp/sidebar-ingestion") / file_id / "derived"


def _load_audio_transcriber() -> Callable[..., dict]:
    global _audio_transcriber
    if _audio_transcriber is not None:
        return _audio_transcriber
    candidate_roots = [
        Path(settings.skills_dir),
        Path(__file__).resolve().parents[2] / "skills",
    ]
    skill_path = None
    for root in candidate_roots:
        candidate = root / "audio-transcribe" / "scripts" / "transcribe_audio.py"
        if candidate.exists():
            skill_path = candidate
            break
    if skill_path is None:
        attempted = ", ".join(str(root / "audio-transcribe" / "scripts" / "transcribe_audio.py") for root in candidate_roots)
        raise IngestionError(
            "TRANSCRIPTION_UNAVAILABLE",
            f"Audio transcription skill not found. Checked: {attempted}",
            retryable=False,
        )
    spec = importlib.util.spec_from_file_location("audio_transcribe_skill", skill_path)
    if not spec or not spec.loader:
        raise IngestionError("TRANSCRIPTION_UNAVAILABLE", "Audio transcription skill could not be loaded", retryable=False)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    transcriber = getattr(module, "transcribe_audio", None)
    if not callable(transcriber):
        raise IngestionError("TRANSCRIPTION_UNAVAILABLE", "Audio transcription entry point missing", retryable=False)
    _audio_transcriber = transcriber
    return transcriber


def _load_youtube_transcriber() -> Callable[..., dict]:
    global _youtube_transcriber
    if _youtube_transcriber is not None:
        return _youtube_transcriber
    candidate_roots = [
        Path(settings.skills_dir),
        Path(__file__).resolve().parents[2] / "skills",
    ]
    skill_path = None
    for root in candidate_roots:
        candidate = root / "youtube-transcribe" / "scripts" / "transcribe_youtube.py"
        if candidate.exists():
            skill_path = candidate
            break
    if skill_path is None:
        attempted = ", ".join(str(root / "youtube-transcribe" / "scripts" / "transcribe_youtube.py") for root in candidate_roots)
        raise IngestionError(
            "VIDEO_TRANSCRIPTION_UNAVAILABLE",
            f"YouTube transcription skill not found. Checked: {attempted}",
            retryable=False,
        )
    spec = importlib.util.spec_from_file_location("youtube_transcribe_skill", skill_path)
    if not spec or not spec.loader:
        raise IngestionError("VIDEO_TRANSCRIPTION_UNAVAILABLE", "YouTube transcription skill could not be loaded", retryable=False)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    transcriber = getattr(module, "transcribe_youtube", None)
    if not callable(transcriber):
        raise IngestionError("VIDEO_TRANSCRIPTION_UNAVAILABLE", "YouTube transcription entry point missing", retryable=False)
    _youtube_transcriber = transcriber
    return transcriber


def _transcribe_youtube(record: IngestedFile) -> tuple[str, dict]:
    if not record.source_url:
        raise IngestionError("INVALID_YOUTUBE_URL", "Missing YouTube URL", retryable=False)
    transcriber = _load_youtube_transcriber()
    temp_root = _derivative_dir(str(record.id))
    temp_root.mkdir(parents=True, exist_ok=True)
    try:
        result = transcriber(
            record.source_url,
            user_id=str(record.user_id),
            output_dir=f"files/{record.id}/ai",
            output_name="ai.md",
            audio_dir="files/videos",
            keep_audio=False,
        )
    except ValueError as exc:
        raise IngestionError("INVALID_YOUTUBE_URL", str(exc), retryable=False) from exc
    except RuntimeError as exc:
        message = str(exc)
        retryable = "rate" in message.lower() or "timeout" in message.lower()
        raise IngestionError("VIDEO_TRANSCRIPTION_FAILED", message, retryable=retryable) from exc
    except Exception as exc:
        raise IngestionError("VIDEO_TRANSCRIPTION_FAILED", "Video transcription failed", retryable=True) from exc

    transcript_path = Path(result.get("transcript_file") or "")
    if not transcript_path.exists():
        raise IngestionError("VIDEO_TRANSCRIPTION_FAILED", "Transcript output missing", retryable=False)
    transcript = transcript_path.read_text(encoding="utf-8", errors="ignore").strip()
    if not transcript:
        transcript = "No transcription available."
    metadata = {
        "provider": "youtube",
        "title": result.get("title"),
        "youtube_url": record.source_url,
        "audio_duration": result.get("audio_duration"),
        "download_duration_seconds": result.get("download_duration_seconds"),
        "transcription_duration_seconds": result.get("transcription_duration_seconds"),
    }
    return transcript, metadata


def _transcribe_audio(source_path: Path, record: IngestedFile) -> str:
    transcriber = _load_audio_transcriber()
    temp_root = _derivative_dir(str(record.id))
    temp_root.mkdir(parents=True, exist_ok=True)
    filename = Path(record.filename_original).name or "audio"
    if "." not in filename:
        extension = _detect_extension(record.filename_original, record.mime_original or "")
        filename = f"{filename}{extension or ''}"
    temp_path = temp_root / filename
    shutil.copyfile(source_path, temp_path)
    try:
        result = transcriber(
            str(temp_path),
            user_id=str(record.user_id),
            output_dir=f"files/{record.id}/ai",
            temp_dir=str(_derivative_dir(str(record.id))),
            response_format="json",
        )
    except RuntimeError as exc:
        raise IngestionError("TRANSCRIPTION_UNAVAILABLE", str(exc), retryable=False) from exc
    except Exception as exc:
        raise IngestionError("TRANSCRIPTION_FAILED", "Audio transcription failed", retryable=True) from exc
    transcript = (result.get("transcript") or "").strip()
    return transcript or "No transcription available."


def _detect_extension(filename: str, mime: str) -> str:
    suffix = Path(filename).suffix
    if suffix:
        return suffix
    if mime == "image/png":
        return ".png"
    if mime in {"image/jpeg", "image/jpg"}:
        return ".jpg"
    return ""


def _guess_text_mime(filename: str) -> str | None:
    ext = Path(filename).suffix.lower()
    if ext in {".md", ".markdown"}:
        return "text/markdown"
    if ext in {".txt", ".log"}:
        return "text/plain"
    if ext == ".csv":
        return "text/csv"
    if ext == ".tsv":
        return "text/tab-separated-values"
    if ext == ".json":
        return "application/json"
    if ext in {".yml", ".yaml"}:
        return "text/plain"
    if ext == ".xls":
        return "application/vnd.ms-excel"
    if ext in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    return None


def _run_libreoffice_convert(source_path: Path, output_dir: Path) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    command = [
        "soffice",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(source_path),
    ]
    try:
        subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=120)
    except FileNotFoundError as exc:
        raise IngestionError("CONVERSION_UNAVAILABLE", "LibreOffice not available", retryable=False) from exc
    except subprocess.TimeoutExpired as exc:
        raise IngestionError("CONVERSION_TIMEOUT", "Conversion timed out", retryable=True) from exc
    except subprocess.CalledProcessError as exc:
        raise IngestionError("CONVERSION_FAILED", "Conversion failed", retryable=False) from exc

    pdf_path = output_dir / (source_path.stem + ".pdf")
    if not pdf_path.exists():
        raise IngestionError("CONVERSION_FAILED", "Converted PDF missing", retryable=False)
    return pdf_path


def _generate_image_thumbnail(image_path: Path, max_size: int = 640) -> bytes | None:
    try:
        with Image.open(image_path) as image:
            image.thumbnail((max_size, max_size))
            output = io.BytesIO()
            image.save(output, format="PNG")
            return output.getvalue()
    except Exception:
        return None


def _generate_pdf_thumbnail(pdf_path: Path, output_dir: Path) -> bytes | None:
    if not shutil.which("pdftoppm"):
        return None
    output_dir.mkdir(parents=True, exist_ok=True)
    output_base = output_dir / "thumb"
    command = [
        "pdftoppm",
        "-f",
        "1",
        "-l",
        "1",
        "-singlefile",
        "-png",
        str(pdf_path),
        str(output_base),
    ]
    try:
        subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=60)
    except (FileNotFoundError, subprocess.TimeoutExpired, subprocess.CalledProcessError):
        return None
    thumb_path = output_dir / "thumb.png"
    if not thumb_path.exists():
        return None
    return thumb_path.read_bytes()


def _merge_split_words(line: str) -> str:
    parts = line.split(" ")
    if len(parts) < 2:
        return line
    merged: list[str] = []
    idx = 0
    while idx < len(parts):
        current = parts[idx]
        if idx + 1 < len(parts):
            next_part = parts[idx + 1]
            if current.isalpha() and next_part.isalpha():
                left = current.lower()
                right = next_part.lower()
                if (
                    (len(current) <= 2 and len(next_part) <= 2)
                    and left not in _COMMON_SHORT_WORDS
                    and right not in _COMMON_SHORT_WORDS
                    and not (current[0].isupper() and next_part[0].isupper())
                ):
                    current = f"{current}{next_part}"
                    idx += 1
        merged.append(current)
        idx += 1
    return " ".join(merged)


def _order_pdf_lines(lines: list[tuple[float, float, str]], page_width: float) -> list[str]:
    if not lines:
        return []
    xs = sorted(x for x, _, _ in lines)
    left_x = xs[0]
    right_x = xs[-1]
    inferred_width = right_x - left_x
    if page_width <= 0:
        page_width = inferred_width
    split_x: float | None = None
    if page_width > 0 and len(xs) >= 6:
        gaps = [(xs[i + 1] - xs[i], i) for i in range(len(xs) - 1)]
        max_gap, gap_index = max(gaps, key=lambda item: item[0])
        if max_gap > page_width * 0.05:
            split_x = (xs[gap_index] + xs[gap_index + 1]) / 2

    def sort_lines(values: list[tuple[float, float, str]]) -> list[str]:
        return [text for _, _, text in sorted(values, key=lambda item: (-item[1], item[0]))]

    if split_x is None:
        return sort_lines(lines)

    left = [line for line in lines if line[0] <= split_x]
    right = [line for line in lines if line[0] > split_x]
    ordered = sort_lines(left)
    if left and right:
        ordered.append("")
    ordered.extend(sort_lines(right))
    return ordered


def _split_line_segments(line: LTTextLine) -> list[tuple[float, str]]:
    chars = [obj for obj in line if isinstance(obj, LTChar)]
    if not chars:
        return []
    widths = [char.width for char in chars if char.width > 0]
    avg_width = statistics.median(widths) if widths else 0.0
    gap_threshold = avg_width * 2.5 if avg_width > 0 else 6.0
    segments: list[tuple[float, str]] = []
    current: list[str] = []
    current_x0: float | None = None
    prev_x1: float | None = None
    for obj in line:
        if isinstance(obj, LTChar):
            if prev_x1 is not None and obj.x0 - prev_x1 > gap_threshold:
                text = "".join(current).strip()
                if text:
                    segments.append((current_x0 or obj.x0, text))
                current = []
                current_x0 = None
            if current_x0 is None:
                current_x0 = obj.x0
            current.append(obj.get_text())
            prev_x1 = obj.x1
        elif isinstance(obj, LTAnno):
            current.append(obj.get_text())
    if current:
        text = "".join(current).strip()
        if text:
            segments.append((current_x0 or chars[0].x0, text))
    return segments


def _should_merge_lines(current: str, next_line: str) -> bool:
    if not next_line:
        return False
    if current.endswith((".", "!", "?", ":", ";")):
        return False
    if current.startswith(("- ", "* ", "• ")):
        return False
    if next_line.startswith(("- ", "* ", "• ")):
        return False
    if next_line[0].islower():
        return True
    return current.endswith(",")


def _join_wrapped_lines(lines: list[str]) -> list[str]:
    merged: list[str] = []
    idx = 0
    while idx < len(lines):
        line = lines[idx]
        if not line:
            merged.append("")
            idx += 1
            continue
        while idx + 1 < len(lines) and _should_merge_lines(line, lines[idx + 1]):
            line = f"{line} {lines[idx + 1]}".strip()
            idx += 1
        merged.append(line)
        idx += 1
    return merged


def _clean_extracted_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = text.replace("\u00ad", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    text = re.sub(r"(\w)-\n(\w)", r"\1-\2", text)
    text = re.sub(r"([a-z]{4,})([A-Z][a-z]{2,})", r"\1 \2", text)
    text = re.sub(r"\b([A-Z][a-z]{2,})([A-Z][a-z]{2,})\b", r"\1 \2", text)
    text = re.sub(r"([a-z])([A-Z]{2,})", r"\1 \2", text)
    text = re.sub(r"(\d)([A-Za-z])", r"\1 \2", text)
    text = re.sub(r"([a-z])(\d)", r"\1 \2", text)
    lines = [line.strip() for line in text.splitlines()]
    lines = [_merge_split_words(line) if line else "" for line in lines]
    lines = _join_wrapped_lines(lines)
    text = "\n".join(lines)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _normalize_table_rows(rows: list[list[object | None]]) -> list[list[str]]:
    if not rows:
        return []
    max_cols = max(len(row) for row in rows)
    normalized = []
    for row in rows:
        padded = [str(cell).strip() if cell is not None else "" for cell in row]
        padded += [""] * (max_cols - len(padded))
        normalized.append(padded)
    last_col = -1
    for col in range(max_cols):
        if any(row[col].strip() for row in normalized):
            last_col = col
    if last_col < 0:
        return []
    trimmed = [row[: last_col + 1] for row in normalized]
    return [row for row in trimmed if any(cell.strip() for cell in row)]


def _table_to_markdown(rows: list[list[object | None]]) -> str | None:
    normalized = _normalize_table_rows(rows)
    if len(normalized) < 2 or len(normalized[0]) < 2:
        return None
    header = normalized[0]
    body = normalized[1:]
    if not body:
        return None
    return tabulate(body, headers=header, tablefmt="github")


def _extract_non_table_text(page: pdfplumber.page.Page, table_bboxes: list[tuple[float, float, float, float]]) -> str:
    words = page.extract_words(x_tolerance=2, y_tolerance=2, keep_blank_chars=False)
    if not words:
        return ""
    filtered = []
    for word in words:
        x0 = float(word["x0"])
        x1 = float(word["x1"])
        top = float(word["top"])
        bottom = float(word["bottom"])
        inside = False
        for bx0, by0, bx1, by1 in table_bboxes:
            if x0 >= bx0 and x1 <= bx1 and top >= by0 and bottom <= by1:
                inside = True
                break
        if not inside:
            filtered.append(word)
    if not filtered:
        return ""
    filtered.sort(key=lambda item: (item["top"], item["x0"]))
    lines: list[list[str]] = []
    current: list[str] = []
    current_top: float | None = None
    for word in filtered:
        top = float(word["top"])
        if current_top is None or abs(top - current_top) <= 3:
            current.append(word["text"])
            current_top = top if current_top is None else current_top
        else:
            lines.append(current)
            current = [word["text"]]
            current_top = top
    if current:
        lines.append(current)
    return "\n".join(" ".join(line).strip() for line in lines if line)


def _extract_pdf_tables(pdf_path: Path) -> tuple[dict[int, list[str]], dict[int, str]]:
    tables_by_page: dict[int, list[str]] = {}
    text_by_page: dict[int, str] = {}
    try:
        with pdfplumber.open(str(pdf_path)) as pdf:
            for index, page in enumerate(pdf.pages):
                tables = page.find_tables()
                if not tables:
                    continue
                table_bboxes = [table.bbox for table in tables]
                markdown_tables = []
                for table in tables:
                    rows = table.extract()
                    markdown = _table_to_markdown(rows or [])
                    if markdown:
                        markdown_tables.append(markdown)
                if markdown_tables:
                    tables_by_page[index] = markdown_tables
                    non_table_text = _extract_non_table_text(page, table_bboxes)
                    if non_table_text:
                        text_by_page[index] = non_table_text
    except Exception as exc:
        logger.warning("pdfplumber table extraction failed (%s).", exc)
    return tables_by_page, text_by_page


def _extract_pdf_text(pdf_path: Path) -> str:
    text = ""
    try:
        tables_by_page, table_text_by_page = _extract_pdf_tables(pdf_path)
        laparams = LAParams(word_margin=0.1, char_margin=2.0, line_margin=0.5)
        sections: list[str] = []
        for index, layout in enumerate(extract_pages(str(pdf_path), laparams=laparams)):
            lines: list[tuple[float, float, str]] = []
            page_width = 0.0
            if hasattr(layout, "bbox"):
                page_width = float(layout.bbox[2] - layout.bbox[0])
            for element in layout:
                if not isinstance(element, LTTextContainer):
                    continue
                for line in element:
                    if not isinstance(line, LTTextLine):
                        continue
                    segments = _split_line_segments(line)
                    if not segments:
                        line_text = line.get_text().strip()
                        if not line_text:
                            continue
                        x0, y0, _, _ = line.bbox
                        lines.append((float(x0), float(y0), line_text))
                        continue
                    _, y0, _, _ = line.bbox
                    for seg_x0, seg_text in segments:
                        if seg_text:
                            lines.append((float(seg_x0), float(y0), seg_text))
            ordered_lines = _order_pdf_lines(lines, page_width)
            page_text = _clean_extracted_text("\n".join(ordered_lines))
            if index in tables_by_page:
                table_text = "\n\n".join(tables_by_page[index])
                non_table_text = table_text_by_page.get(index, "")
                page_text = _clean_extracted_text(non_table_text) if non_table_text else ""
                if page_text:
                    page_text = f"{page_text}\n\n{table_text}"
                else:
                    page_text = table_text
            sections.append(f"## Page {index + 1}\n\n{page_text.strip()}")
        text = "\n\n".join(sections).strip()
    except Exception as exc:
        logger.warning("PDFMiner extraction failed (%s). Falling back to pypdf.", exc)
        try:
            reader = PdfReader(str(pdf_path))
            sections = []
            for index, page in enumerate(reader.pages, start=1):
                page_text = page.extract_text() or ""
                page_text = _clean_extracted_text(page_text)
                sections.append(f"## Page {index}\n\n{page_text.strip()}")
            text = "\n\n".join(sections).strip()
        except Exception as fallback_exc:
            logger.warning("pypdf extraction failed (%s).", fallback_exc)
            text = ""
    return _clean_extracted_text(text)


def _extract_pptx_text(pptx_path: Path) -> str:
    presentation = Presentation(str(pptx_path))
    sections = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        content = "\n".join(texts).strip()
        sections.append(f"## Slide {index}\n\n{content}")
    return "\n\n".join(sections).strip()


def _extract_docx_text(docx_path: Path) -> str:
    document = Document(str(docx_path))
    parts: list[str] = []
    for paragraph in document.paragraphs:
        if paragraph.text:
            parts.append(paragraph.text)
    for table in document.tables:
        rows = []
        for row in table.rows:
            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        if rows:
            parts.append("\n".join(rows))
    return "\n\n".join(parts).strip()


def _stringify_cell(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    if isinstance(value, datetime):
        return value.isoformat()
    if isinstance(value, date):
        return value.isoformat()
    return str(value)


def _column_index(column: str) -> int:
    index = 0
    for char in column.upper():
        if not ("A" <= char <= "Z"):
            break
        index = index * 26 + (ord(char) - ord("A") + 1)
    return max(index - 1, 0)


def _parse_cell_reference(cell_ref: str) -> tuple[int, int]:
    if not isinstance(cell_ref, str):
        return 0, 0
    column_part = ""
    row_part = ""
    for char in cell_ref:
        if char.isdigit():
            row_part += char
        else:
            column_part += char
    if not row_part:
        return 0, 0
    return int(row_part) - 1, _column_index(column_part)


def _parse_range(
    range_ref: object, row_count: int | None = None, col_count: int | None = None
) -> tuple[int, int, int, int]:
    if isinstance(range_ref, str):
        if ":" in range_ref:
            start_ref, end_ref = range_ref.split(":", 1)
        else:
            start_ref = end_ref = range_ref
        start_row, start_col = _parse_cell_reference(start_ref)
        end_row, end_col = _parse_cell_reference(end_ref)
        return (
            min(start_row, end_row),
            min(start_col, end_col),
            max(start_row, end_row),
            max(start_col, end_col),
        )
    if isinstance(range_ref, (list, tuple)):
        parts = list(range_ref)
        numeric_candidates: list[tuple[int, int, int, int]] = []
        if len(parts) == 2:
            start_ref, end_ref = parts
            if isinstance(start_ref, str) and isinstance(end_ref, str):
                start_row, start_col = _parse_cell_reference(start_ref)
                end_row, end_col = _parse_cell_reference(end_ref)
                return (
                    min(start_row, end_row),
                    min(start_col, end_col),
                    max(start_row, end_row),
                    max(start_col, end_col),
                )
            if (
                isinstance(start_ref, (list, tuple))
                and isinstance(end_ref, (list, tuple))
                and len(start_ref) >= 2
                and len(end_ref) >= 2
            ):
                start_row, start_col = int(start_ref[0]), int(start_ref[1])
                end_row, end_col = int(end_ref[0]), int(end_ref[1])
                numeric_candidates.append(
                    (
                        min(start_row, end_row),
                        min(start_col, end_col),
                        max(start_row, end_row),
                        max(start_col, end_col),
                    )
                )
        if len(parts) == 4 and all(isinstance(part, (int, float)) for part in parts):
            start_row, start_col, end_row, end_col = [int(part) for part in parts]
            numeric_candidates.append(
                (
                    min(start_row, end_row),
                    min(start_col, end_col),
                    max(start_row, end_row),
                    max(start_col, end_col),
                )
            )
            numeric_candidates.append(
                (
                    min(start_col, end_col),
                    min(start_row, end_row),
                    max(start_col, end_col),
                    max(start_row, end_row),
                )
            )
        if numeric_candidates:
            for candidate in numeric_candidates:
                if row_count is None or col_count is None:
                    return candidate
                if candidate[0] < row_count and candidate[1] < col_count:
                    return candidate
            adjusted_candidates: list[tuple[int, int, int, int]] = []
            for start_row, start_col, end_row, end_col in numeric_candidates:
                if min(start_row, start_col, end_row, end_col) >= 1:
                    adjusted_candidates.append(
                        (start_row - 1, start_col - 1, end_row - 1, end_col - 1)
                    )
            for candidate in adjusted_candidates:
                if row_count is None or col_count is None:
                    return candidate
                if candidate[0] < row_count and candidate[1] < col_count:
                    return candidate
    return 0, 0, 0, 0


def _normalize_grid(rows: list[list[object]]) -> list[list[object]]:
    if not rows:
        return []
    max_columns = max(len(row) for row in rows)
    normalized = []
    for row in rows:
        padded = list(row) + [None] * (max_columns - len(row))
        normalized.append(padded)
    return normalized


def _expand_merged_cells(rows: list[list[object]], merged_ranges: list[object]) -> list[list[object]]:
    if not rows or not merged_ranges:
        return rows
    max_columns = max(len(row) for row in rows) if rows else 0
    row_count = len(rows)
    for range_ref in merged_ranges:
        start_row, start_col, end_row, end_col = _parse_range(range_ref, row_count, max_columns)
        required_rows = end_row + 1
        required_cols = end_col + 1
        while len(rows) < required_rows:
            rows.append([None] * max_columns)
        if required_cols > max_columns:
            for row in rows:
                row.extend([None] * (required_cols - len(row)))
            max_columns = required_cols

        value = None
        for row_index in range(start_row, end_row + 1):
            for col_index in range(start_col, end_col + 1):
                cell_value = rows[row_index][col_index]
                if cell_value not in (None, ""):
                    value = cell_value
                    break
            if value is not None:
                break
        if value is None:
            continue
        for row_index in range(start_row, end_row + 1):
            for col_index in range(start_col, end_col + 1):
                rows[row_index][col_index] = value
    return rows


def _build_sheet_payload(name: str, rows: list[list[str]]) -> dict:
    return {
        "name": name,
        "rows": rows,
        "header_row": 0 if rows else None,
    }


def _read_csv_rows(content: bytes, delimiter: str) -> list[list[str]]:
    text = content.decode("utf-8-sig", errors="ignore")
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    rows: list[list[str]] = []
    for row in reader:
        rows.append([_stringify_cell(cell) for cell in row])
    return rows


def _extract_spreadsheet_data(source_path: Path, mime: str, filename: str) -> dict:
    content = source_path.read_bytes()
    lower_name = filename.lower()
    if mime in {"text/csv", "application/csv"} or lower_name.endswith(".csv"):
        rows = _read_csv_rows(content, ",")
        return {"sheets": [_build_sheet_payload(Path(filename).stem or "Sheet1", rows)]}
    if mime in {"text/tab-separated-values", "text/tsv"} or lower_name.endswith(".tsv"):
        rows = _read_csv_rows(content, "\t")
        return {"sheets": [_build_sheet_payload(Path(filename).stem or "Sheet1", rows)]}

    try:
        workbook = CalamineWorkbook.from_path(str(source_path))
    except (CalamineError, ZipError, XmlError) as exc:
        raise IngestionError(
            "INVALID_XLSX",
            "Invalid spreadsheet file. Re-save it as .xlsx in Excel or Google Sheets.",
            retryable=False,
        ) from exc

    sheets = []
    try:
        for sheet_name in workbook.sheet_names:
            sheet = workbook.get_sheet_by_name(sheet_name)
            raw_rows = sheet.to_python(skip_empty_area=False) or []
            normalized = _normalize_grid(raw_rows)
            merged_ranges = getattr(sheet, "merged_cell_ranges", None) or []
            normalized = _expand_merged_cells(normalized, list(merged_ranges))
            rows = [[_stringify_cell(cell) for cell in row] for row in normalized]
            sheets.append(_build_sheet_payload(sheet_name, rows))
    finally:
        workbook.close()
    return {"sheets": sheets}


def _build_derivatives(record: IngestedFile, source_path: Path) -> list[DerivativePayload]:
    mime = (record.mime_original or "application/octet-stream").split(";")[0].strip().lower()
    if not mime:
        mime = "application/octet-stream"
    if mime == "application/octet-stream":
        guessed_mime = _guess_text_mime(record.filename_original)
        if guessed_mime:
            mime = guessed_mime
    file_id = str(record.id)
    user_prefix = f"{record.user_id}/files/{file_id}"
    content = source_path.read_bytes()
    extraction_text = ""
    viewer_payload: DerivativePayload | None = None
    thumb_payload: DerivativePayload | None = None
    thumb_bytes: bytes | None = None

    if mime == "application/pdf":
        viewer_payload = DerivativePayload(
            kind="viewer_pdf",
            storage_key=f"{user_prefix}/derivatives/viewer.pdf",
            mime="application/pdf",
            size_bytes=len(content),
            sha256=sha256(content).hexdigest(),
            content=content,
        )
        extraction_text = _extract_pdf_text(source_path)
        thumb_bytes = _generate_pdf_thumbnail(source_path, _derivative_dir(file_id))
    elif mime.startswith("image/"):
        extension = _detect_extension(record.filename_original, mime)
        viewer_payload = DerivativePayload(
            kind="image_original",
            storage_key=f"{user_prefix}/derivatives/image{extension or ''}",
            mime=mime,
            size_bytes=len(content),
            sha256=sha256(content).hexdigest(),
            content=content,
        )
        extraction_text = "Image file. No text extraction available."
        thumb_bytes = _generate_image_thumbnail(source_path)
    elif mime.startswith("audio/"):
        extension = _detect_extension(record.filename_original, mime)
        viewer_payload = DerivativePayload(
            kind="audio_original",
            storage_key=f"{user_prefix}/derivatives/audio{extension or ''}",
            mime=mime,
            size_bytes=len(content),
            sha256=sha256(content).hexdigest(),
            content=content,
        )
        extraction_text = _transcribe_audio(source_path, record)
    elif mime in {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }:
        pdf_path = _run_libreoffice_convert(source_path, _derivative_dir(file_id))
        pdf_bytes = pdf_path.read_bytes()
        viewer_payload = DerivativePayload(
            kind="viewer_pdf",
            storage_key=f"{user_prefix}/derivatives/viewer.pdf",
            mime="application/pdf",
            size_bytes=len(pdf_bytes),
            sha256=sha256(pdf_bytes).hexdigest(),
            content=pdf_bytes,
        )
        if mime.endswith("wordprocessingml.document"):
            extraction_text = _extract_docx_text(source_path)
        else:
            extraction_text = _extract_pptx_text(source_path)
        thumb_bytes = _generate_pdf_thumbnail(pdf_path, _derivative_dir(file_id))
    elif (
        mime
        in {
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
            "text/csv",
            "application/csv",
            "text/tab-separated-values",
            "text/tsv",
        }
        or record.filename_original.lower().endswith((".csv", ".tsv"))
    ):
        data = _extract_spreadsheet_data(source_path, mime, record.filename_original)
        data_json = json.dumps(data, ensure_ascii=False)
        data_bytes = data_json.encode("utf-8")
        viewer_payload = DerivativePayload(
            kind="viewer_json",
            storage_key=f"{user_prefix}/derivatives/data.json",
            mime="application/json",
            size_bytes=len(data_bytes),
            sha256=sha256(data_bytes).hexdigest(),
            content=data_bytes,
        )
        extraction_text = data_json
    elif mime.startswith("text/") or mime == "application/json":
        viewer_payload = DerivativePayload(
            kind="text_original",
            storage_key=f"{user_prefix}/derivatives/source.txt",
            mime=mime,
            size_bytes=len(content),
            sha256=sha256(content).hexdigest(),
            content=content,
        )
        extraction_text = content.decode("utf-8", errors="ignore").strip()
    else:
        raise IngestionError("UNSUPPORTED_TYPE", "Unsupported file type", retryable=False)

    if viewer_payload is None:
        raise IngestionError("DERIVATIVE_MISSING", "Viewer asset missing", retryable=False)

    if thumb_bytes:
        thumb_payload = DerivativePayload(
            kind="thumb_png",
            storage_key=f"{user_prefix}/derivatives/thumb.png",
            mime="image/png",
            size_bytes=len(thumb_bytes),
            sha256=sha256(thumb_bytes).hexdigest(),
            content=thumb_bytes,
        )

    ai_body = (
        "---\n"
        f"file_id: {record.id}\n"
        f"source_filename: {record.filename_original}\n"
        f"source_mime: {record.mime_original}\n"
        f"created_at: {record.created_at.isoformat()}\n"
        f"sha256: {record.sha256}\n"
        "derivatives:\n"
        f"  {viewer_payload.kind}: true\n"
        "---\n\n"
        f"{extraction_text or 'No text extraction available.'}"
    )
    ai_bytes = ai_body.encode("utf-8")
    ai_md = DerivativePayload(
        kind="ai_md",
        storage_key=f"{user_prefix}/ai/ai.md",
        mime="text/markdown",
        size_bytes=len(ai_bytes),
        sha256=sha256(ai_bytes).hexdigest(),
        content=ai_bytes,
    )

    derivatives = [viewer_payload, ai_md]
    if thumb_payload:
        derivatives.append(thumb_payload)
    return derivatives


def _build_youtube_derivatives(
    record: IngestedFile,
    transcript: str,
    metadata: dict,
) -> list[DerivativePayload]:
    user_prefix = f"{record.user_id}/files/{record.id}"
    title = metadata.get("title") or record.filename_original
    ai_body = (
        "---\n"
        f"file_id: {record.id}\n"
        f"source_title: {title}\n"
        f"source_url: {record.source_url}\n"
        f"source_mime: {record.mime_original}\n"
        f"created_at: {record.created_at.isoformat()}\n"
        "derivatives:\n"
        "  ai_md: true\n"
        "---\n\n"
        f"{transcript or 'No transcription available.'}"
    )
    ai_bytes = ai_body.encode("utf-8")
    ai_md = DerivativePayload(
        kind="ai_md",
        storage_key=f"{user_prefix}/ai/ai.md",
        mime="text/markdown",
        size_bytes=len(ai_bytes),
        sha256=sha256(ai_bytes).hexdigest(),
        content=ai_bytes,
    )
    return [ai_md]


def _process_youtube_job(db, job: FileProcessingJob, record: IngestedFile) -> None:
    logger.info("Processing YouTube file_id=%s url=%s", record.id, record.source_url)
    for stage in PIPELINE_STAGES:
        db.refresh(job)
        if job.status in {"paused", "canceled"}:
            raise IngestionError("JOB_HALTED", "Job halted by user", retryable=False)
        _set_stage(db, job, stage)
        logger.info("Stage %s file_id=%s", stage, record.id)
        _refresh_lease(db, job)

        if stage == "validating":
            if not record.source_url:
                raise IngestionError("INVALID_YOUTUBE_URL", "Missing YouTube URL", retryable=False)
        elif stage in {"converting", "extracting", "ai_md", "thumb"}:
            time.sleep(0.05)
        elif stage == "finalizing":
            transcript, metadata = _transcribe_youtube(record)
            existing_meta = record.source_metadata or {}
            metadata = {**existing_meta, **metadata}
            metadata.setdefault("provider", "youtube")
            metadata.setdefault("youtube_url", record.source_url)
            title = metadata.get("title")
            if title:
                record.filename_original = title
            record.source_metadata = metadata
            db.commit()

            derivatives = _build_youtube_derivatives(record, transcript, metadata)
            storage = get_storage_backend()
            for item in derivatives:
                storage.put_object(item.storage_key, item.content, content_type=item.mime)

            db.query(FileDerivative).filter(FileDerivative.file_id == record.id).delete()
            now = _now()
            for item in derivatives:
                db.add(
                    FileDerivative(
                        file_id=record.id,
                        kind=item.kind,
                        storage_key=item.storage_key,
                        mime=item.mime,
                        size_bytes=item.size_bytes,
                        sha256=item.sha256,
                        created_at=now,
                    )
                )
            db.commit()


def worker_loop() -> None:
    worker_id = os.getenv("INGESTION_WORKER_ID") or f"worker-{uuid4()}"
    worker_user_id = os.getenv("INGESTION_WORKER_USER_ID") or settings.default_user_id
    if shutil.which("soffice") is None:
        logger.warning(
            "LibreOffice (soffice) not found. DOCX/XLSX/PPTX conversion will fail."
        )
    db_url = settings.database_url
    if db_url:
        safe_url = db_url
        if "://" in safe_url:
            scheme, rest = safe_url.split("://", 1)
            if "@" in rest:
                _, rest = rest.split("@", 1)
            safe_url = f"{scheme}://{rest}"
        logger.info("Ingestion DB URL: %s", safe_url)
    while True:
        with SessionLocal() as db:
            if worker_user_id:
                db.execute(text("SET app.user_id = :user_id"), {"user_id": worker_user_id})
            _requeue_stalled_jobs(db)
            job = _claim_job(db, worker_id)
            if not job:
                time.sleep(SLEEP_SECONDS)
                continue

            try:
                record = _get_file(db, job)
                if record.user_id:
                    db.execute(text("SET app.user_id = :user_id"), {"user_id": str(record.user_id)})
                if record.source_url:
                    _process_youtube_job(db, job, record)
                    db.refresh(job)
                    if job.status not in {"paused", "canceled"}:
                        _mark_ready(db, job)
                    continue

                source_path = _staging_path(str(record.id))
                if not source_path.exists():
                    raise IngestionError("SOURCE_MISSING", "Uploaded file not found", retryable=False)

                logger.info("Processing file_id=%s name=%s", record.id, record.filename_original)
                for stage in PIPELINE_STAGES:
                    db.refresh(job)
                    if job.status in {"paused", "canceled"}:
                        raise IngestionError("JOB_HALTED", "Job halted by user", retryable=False)
                    _set_stage(db, job, stage)
                    logger.info("Stage %s file_id=%s", stage, record.id)
                    _refresh_lease(db, job)

                    if stage == "validating":
                        if record.size_bytes <= 0:
                            raise IngestionError("FILE_EMPTY", "Uploaded file is empty", retryable=False)
                    elif stage in {"converting", "extracting", "ai_md", "thumb"}:
                        time.sleep(0.05)
                    elif stage == "finalizing":
                        derivatives = _build_derivatives(record, source_path)
                        storage = get_storage_backend()
                        for item in derivatives:
                            storage.put_object(item.storage_key, item.content, content_type=item.mime)

                        db.query(FileDerivative).filter(FileDerivative.file_id == record.id).delete()
                        now = _now()
                        for item in derivatives:
                            db.add(
                                FileDerivative(
                                    file_id=record.id,
                                    kind=item.kind,
                                    storage_key=item.storage_key,
                                    mime=item.mime,
                                    size_bytes=item.size_bytes,
                                    sha256=item.sha256,
                                    created_at=now,
                                )
                            )
                        db.commit()

                db.refresh(job)
                if job.status not in {"paused", "canceled"}:
                    _mark_ready(db, job)
                    _cleanup_staging(str(record.id))
            except IngestionError as error:
                if error.code == "JOB_HALTED":
                    continue
                db.rollback()
                _retry_or_fail(db, job, error)
                if job.status == "failed" and _should_cleanup_after_failure(error):
                    _cleanup_staging(str(job.file_id))
            except Exception as error:
                unknown_error = IngestionError("UNKNOWN_ERROR", str(error), retryable=True)
                db.rollback()
                _retry_or_fail(db, job, unknown_error)
                if job.status == "failed" and _should_cleanup_after_failure(unknown_error):
                    _cleanup_staging(str(job.file_id))
        time.sleep(SLEEP_SECONDS)


if __name__ == "__main__":
    while True:
        try:
            worker_loop()
        except Exception:
            logger.exception("Ingestion worker crashed, restarting after delay")
            time.sleep(2)
